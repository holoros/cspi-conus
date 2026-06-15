#!/usr/bin/env python3
"""Build CSPI v0.10 collaborators briefing deck (Committee Update mode).

Audience: CSPI scientific collaborators, FIA Carbon Field Pilot leads, CFRU
partners, NSRC collaborators. Target read time: roughly 12 minutes.

Branding: CRSF primary green #1A3D28, Aptos Display titles, Aptos body.
Mode: Committee Update / Brief. White background, takeaway band optional,
8-20 words per slide, no speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUT = "/sessions/quirky-peaceful-babbage/mnt/outputs/CSPI_v10_collaborators_overview.pptx"

# CRSF brand
CRSF_GREEN = RGBColor(0x1A, 0x3D, 0x28)
CRSF_LIGHT = RGBColor(0x6F, 0x8C, 0x5A)
CHARCOAL   = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Component-themed accent colors (matching F11 decision tree)
COLOR_ESI  = RGBColor(0x5B, 0x68, 0xB8)
COLOR_BGI  = RGBColor(0xE1, 0x55, 0x54)
COLOR_ASYM = RGBColor(0x6F, 0x8C, 0x5A)
COLOR_NPP  = RGBColor(0xC6, 0x8C, 0x42)
COLOR_CSPI = RGBColor(0x8B, 0x5D, 0xA6)

CRSF_LOGO = "/sessions/quirky-peaceful-babbage/mnt/aweiskittel/Documents/Claude/CRSF-Cowork/Logos_CRSFgraphics/Center Logos/CRSF/CRSF logo (2018)/1. Full Color/CRSF logo.png"
F8 = "/sessions/quirky-peaceful-babbage/mnt/bgi-cspi-conus/v5/figures_v10/F8_siteclcd_by_measure.png"
F9 = "/sessions/quirky-peaceful-babbage/mnt/bgi-cspi-conus/v5/figures_v10/F9_stand_age_correlation.png"
F10 = "/sessions/quirky-peaceful-babbage/mnt/bgi-cspi-conus/v5/figures_v10/F10_per_species_ESI_BGI.png"
F11 = "/sessions/quirky-peaceful-babbage/mnt/bgi-cspi-conus/v5/figures_v10/F11_decision_tree.png"
F7 = "/sessions/quirky-peaceful-babbage/mnt/bgi-cspi-conus/v5/figures_v10/F7_cspi_v21_3c_map.png"
F1 = "/sessions/quirky-peaceful-babbage/mnt/bgi-cspi-conus/v5/figures_v04/F1_cor_heatmap.png"

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=CHARCOAL,
             font="Aptos", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def add_band(slide, text, color=CRSF_GREEN):
    """Takeaway band along the bottom."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.55),
                                Inches(13.333), Inches(0.95))
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    tf = bg.text_frame
    tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.name = "Aptos"; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = WHITE

def add_logo(slide, x=0.4, y=0.35, size=1.0):
    try:
        slide.shapes.add_picture(CRSF_LOGO, Inches(x), Inches(y),
                                 width=Inches(size), height=Inches(size))
    except Exception:
        pass

def add_rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    return s

def add_image(slide, path, x, y, w, h):
    try:
        return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                        width=Inches(w), height=Inches(h))
    except Exception as e:
        # Placeholder rect
        rect = add_rect(slide, x, y, w, h, RGBColor(0xEE, 0xEE, 0xEE))
        add_text(slide, x, y+h/2-0.3, w, 0.6,
                 f"[image: {Path(path).name}]",
                 size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)
        return rect

# ----------------------------------------------------------------------------
# Slide 1: Title
# ----------------------------------------------------------------------------
s = blank()
# Dark green wash on title
bg = add_rect(s, 0, 0, 13.333, 7.5, CRSF_GREEN)
# Logo top right
try:
    s.shapes.add_picture(CRSF_LOGO, Inches(11.6), Inches(0.4),
                         width=Inches(1.3), height=Inches(1.3))
except Exception:
    pass

add_text(s, 0.8, 1.5, 12.0, 1.6,
         "Beyond site index",
         size=60, bold=True, color=WHITE, font="Aptos Display")
add_text(s, 0.8, 2.9, 12.0, 1.6,
         "FIA's classification disagrees with site index but agrees with a composite measure of forest productivity",
         size=24, color=WHITE, font="Aptos")
add_text(s, 0.8, 5.0, 12.0, 0.5,
         "Aaron R. Weiskittel",
         size=22, color=WHITE, bold=True)
add_text(s, 0.8, 5.5, 12.0, 0.4,
         "University of Maine | Center for Research on Sustainable Forests",
         size=16, color=WHITE)
add_text(s, 0.8, 5.95, 12.0, 0.4,
         "Collaborators briefing | 13 June 2026 | CSPI v2.0.0 release",
         size=14, color=CRSF_LIGHT)
add_text(s, 0.8, 6.85, 12.0, 0.4,
         "Data concept 10.5281/zenodo.20515034 | Analytical chain 10.5281/zenodo.20693106 | v2.0.0 10.5281/zenodo.20663652",
         size=12, color=CRSF_LIGHT)

# ----------------------------------------------------------------------------
# Slide 2: Why this matters now
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "What changed since v1.0.0", size=32, bold=True,
         color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.45,
         "The release line shifted from a single-metric predicted site index to a multi-metric composite.",
         size=18, color=MED_GRAY)

# Three cards: v1, v2, v0.10 reframe
cards = [
    ("v1.0.0 (June 2026)",
     "Single-metric CSPI v3 and v4 at 30 m and 1 km",
     "Predicted site index from FIA SICOND target",
     "Still available at 10.5281/zenodo.20515035",
     COLOR_ESI),
    ("v2.0.0 (June 2026)",
     "Multi-metric composite of ESI + BGI + Asym",
     "Unified North America height-and-age target",
     "Each component ships as a standalone layer",
     COLOR_CSPI),
    ("v0.10 manuscript reframe",
     "Site productivity is empirically multi-dimensional",
     "Single-metric studies miss species, age, regime structure",
     "FIA SITECLCD tracks biomass growth, not site index",
     COLOR_BGI),
]

for i, (title, l1, l2, l3, col) in enumerate(cards):
    cx = 0.7 + i * 4.25
    # left-edge accent bar
    add_rect(s, cx, 1.85, 0.08, 4.6, col)
    # card body
    body = add_rect(s, cx + 0.12, 1.85, 3.95, 4.6, RGBColor(0xF7, 0xF7, 0xF5))
    add_text(s, cx + 0.3, 2.0, 3.7, 0.6, title, size=20, bold=True, color=col,
             font="Aptos Display")
    add_text(s, cx + 0.3, 2.7, 3.7, 1.0, l1, size=15, color=CHARCOAL)
    add_text(s, cx + 0.3, 3.9, 3.7, 1.0, l2, size=15, color=CHARCOAL)
    add_text(s, cx + 0.3, 5.1, 3.7, 1.2, l3, size=15, color=CHARCOAL)

add_band(s, "The release is the composite. The conceptual paper says productivity is multi-dimensional.")

# ----------------------------------------------------------------------------
# Slide 3: Headline metric callout
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "Headline empirical finding", size=32, bold=True,
         color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "FIA's own seven-class productivity rating (SITECLCD) tracks biomass growth, not site index.",
         size=18, color=MED_GRAY)

# Large numbers row
def big_num(x, value, label, sublabel, color):
    add_text(s, x, 2.4, 3.0, 1.4, value, size=60, bold=True, color=color,
             font="Aptos Display", align=PP_ALIGN.CENTER)
    add_text(s, x, 4.0, 3.0, 0.5, label, size=18, bold=True, color=CHARCOAL,
             align=PP_ALIGN.CENTER)
    add_text(s, x, 4.5, 3.0, 0.8, sublabel, size=14, color=MED_GRAY,
             align=PP_ALIGN.CENTER)

big_num(1.0, "0.808", "BGI alone", "OOB R² predicting FIA SITECLCD", COLOR_BGI)
big_num(5.2, "0.751", "ESI alone", "OOB R² predicting FIA SITECLCD", COLOR_ESI)
big_num(9.4, "+0.06", "BGI advantage", "Six percentage-point gap", CRSF_GREEN)

add_text(s, 1.0, 5.6, 11.4, 0.7,
         "Random forest at 63,310 FIA plots. NPP alone reaches only R² = 0.13.",
         size=15, color=MED_GRAY, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.1, 11.4, 0.5,
         "FIA's classification has been measuring biomass growth, not height growth potential. For seven decades.",
         size=16, color=CHARCOAL, align=PP_ALIGN.CENTER, bold=True)

add_band(s, "Single-metric site index studies have been carrying biomass-growth structure unlabeled.")

# ----------------------------------------------------------------------------
# Slide 4: SITECLCD breakdown with F8 figure
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "Mean productivity by FIA site class", size=32, bold=True,
         color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "ESI is flat across the seven classes. BGI drops by a factor of 2.3. The composite tracks the classification.",
         size=16, color=MED_GRAY)
add_image(s, F8, 1.5, 1.85, 10.3, 4.7)
add_band(s, "If site index were what FIA SITECLCD measures, the ESI line would tilt with the classes.")

# ----------------------------------------------------------------------------
# Slide 5: SICOND vs ESI orthogonality
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "Hierarchical ecoregion weighting boosts SITECLCD recovery",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "Equal-weight + ecoregion composite reaches R² = 0.832 vs 0.812 alone (CI non-overlapping).",
         size=15, color=MED_GRAY)

# Two callout numbers
def callout_pair(x, y, val, line1, line2, color):
    add_text(s, x, y, 4.0, 1.1, val, size=50, bold=True, color=color,
             font="Aptos Display", align=PP_ALIGN.CENTER)
    add_text(s, x, y+1.2, 4.0, 0.5, line1, size=16, bold=True, color=CHARCOAL,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y+1.7, 4.0, 0.5, line2, size=13, color=MED_GRAY,
             align=PP_ALIGN.CENTER)

callout_pair(1.5, 2.0, "R² 0.832", "equal + hier composite",
             "Bootstrap 95% CI [0.829, 0.835]; +0.020 over equal alone", COLOR_BGI)
callout_pair(8.0, 2.0, "50 L3", "EPA Level III ecoregions",
             "Mean w_ESI=1.77, w_BGI=0.21, w_Asym=0.34 across regions", COLOR_ESI)

add_text(s, 1.5, 4.5, 10.3, 0.55,
         "Per-region PC1 captures 48% to 84% of variance; Western ESI-dominated, Eastern Asym-dominated.",
         size=17, color=CHARCOAL, bold=True)
add_text(s, 1.5, 5.15, 10.3, 1.35,
         "When you swap SICOND for the unified-target ESI in growth-and-yield work, expect the calibrations to shift. "
         "The shift will be in the direction of less biomass-growth information and more pure height-growth-potential information. "
         "Whether that shift is desirable depends on your application.",
         size=15, color=CHARCOAL)

add_band(s, "Multilevel ecoregion: +0.020 R² gain (CI non-overlapping); SITECLCD result remains the headline anchor.")

# ----------------------------------------------------------------------------
# Slide 6: Stand-age sign flip with F9
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "The relationship flips at 120 years",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "ESI–BGI correlation goes from −0.57 in middle-age stands to +0.33 in old growth.",
         size=15, color=MED_GRAY)
add_image(s, F9, 1.5, 1.85, 10.3, 4.7)
add_band(s, "A productivity surface trained pooled across stand ages is wrong in opposite directions in young vs old stands.")

# ----------------------------------------------------------------------------
# Slide 7: Per-species heterogeneity with F10
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "Species range from 0.028 to 0.492",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "Yellow-poplar essentially independent (r = 0.028); Douglas-fir strongly positive (r = 0.492).",
         size=15, color=MED_GRAY)
add_image(s, F10, 1.5, 1.85, 10.3, 4.7)
add_band(s, "A national productivity ranking implicitly assumes one ESI–BGI relationship. Six species say otherwise.")

# ----------------------------------------------------------------------------
# Slide 8: Decision tree (F11)
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "Which productivity measure should you use?",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "Five operational measures, five primary applications. Pick by what you actually want to know.",
         size=15, color=MED_GRAY)
add_image(s, F11, 1.0, 1.85, 11.3, 4.7)
add_band(s, "If your application is not a clean match for one measure, the composite is the recommended default.")

# ----------------------------------------------------------------------------
# Slide 9: The 30 m spatial map
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "CSPI v2 spatial pattern over CONUS",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "30 m grid, z-score equal-weight average of ESI v7 + BGI + Asym, rescaled to 0–100.",
         size=15, color=MED_GRAY)
add_image(s, F7, 1.5, 1.85, 10.3, 4.7)
add_band(s, "Pacific Northwest coastal forest band highest; southwestern desert margin lowest; central hardwoods intermediate.")

# ----------------------------------------------------------------------------
# Slide 10: What is in the deposit
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "What you get when you cite the v2 deposit",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "Zenodo concept 10.5281/zenodo.20515034 always resolves to the latest version. v2.0.0 lightweight files are live now.",
         size=14, color=MED_GRAY)

items = [
    ("ESI v6", "Recommended 1 km predicted site index", "38 MB", COLOR_ESI),
    ("CSPI v2", "Multi-metric composite metadata + dictionary", "30 KB", COLOR_CSPI),
    ("MOD17 NPP mean", "Annual NPP 2017–2023, 500 m", "130 MB", COLOR_NPP),
    ("MOD17 NPP CV", "Interannual stability layer", "149 MB", COLOR_NPP),
    ("Documentation", "README, dictionary, NEWS, CITATION.cff", "few MB", CRSF_GREEN),
    ("v2.1.0 follow-up", "30 m surfaces, awaiting Zenodo quota", "pending", LIGHT_GRAY),
]
y = 1.95
for (name, desc, size_str, col) in items:
    add_rect(s, 1.5, y, 0.08, 0.65, col)
    add_text(s, 1.7, y+0.05, 2.5, 0.55, name, size=17, bold=True, color=col,
             font="Aptos Display")
    add_text(s, 4.3, y+0.05, 6.0, 0.55, desc, size=14, color=CHARCOAL)
    add_text(s, 10.3, y+0.05, 1.6, 0.55, size_str, size=14, color=MED_GRAY,
             align=PP_ALIGN.RIGHT)
    y += 0.74

add_text(s, 1.5, 6.3, 10.3, 0.25,
         "v1.0.0 SICOND-target surfaces remain available at 10.5281/zenodo.20515035 for legacy use.",
         size=12, color=LIGHT_GRAY)

add_band(s, "Component layers ship separately so you can decompose the composite or use the single measure you need.")

# ----------------------------------------------------------------------------
# Slide 11: Citation guidance
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "How to cite for your work",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "Cite the concept DOI for stability across versions; the version DOI for a specific snapshot.",
         size=14, color=MED_GRAY)

# Two cards: "If you used..." and "Cite as..."
def cite_card(x, y, w, h, kind, text, col):
    add_rect(s, x, y, 0.08, h, col)
    add_rect(s, x+0.12, y, w-0.12, h, RGBColor(0xF7, 0xF7, 0xF5))
    add_text(s, x+0.3, y+0.15, w-0.42, 0.5, kind, size=15, bold=True, color=col,
             font="Aptos Display")
    add_text(s, x+0.3, y+0.7, w-0.42, h-0.85, text, size=12, color=CHARCOAL,
             font="Aptos")

cite_card(0.8, 2.0, 6.0, 4.4, "If you used the composite",
          "Weiskittel, A.R. (2026). Composite Site Productivity Index (CSPI) v2.0.0 "
          "for the conterminous United States. Zenodo.\n"
          "https://doi.org/10.5281/zenodo.20663652\n\n"
          "Or, for a version-stable reference:\n"
          "Weiskittel, A.R. Composite Site Productivity Index (CSPI). Zenodo concept.\n"
          "https://doi.org/10.5281/zenodo.20515034\n\n"
          "Also cite the v0.10 multi-dimensional paper once it is in press at FEM.",
          COLOR_CSPI)

cite_card(6.7, 2.0, 6.0, 4.4, "If you used a single component",
          "Cite the same record but specify the band you used:\n"
          "...ESI v7 layer from CSPI v2.0.0\n"
          "...BGI 30 m layer from CSPI v2.0.0\n"
          "...Asym 30 m layer from CSPI v2.0.0\n"
          "...MOD17 NPP 500 m layer from CSPI v2.0.0\n\n"
          "For SICOND, cite the FIA Phase 2 condition table directly, not this release.",
          COLOR_BGI)

add_band(s, "Concept DOI for stability; version DOI for snapshot; the manuscript for the conceptual reframe.")

# ----------------------------------------------------------------------------
# Slide 12: What is next and how to engage
# ----------------------------------------------------------------------------
s = blank()
add_logo(s)
add_text(s, 1.7, 0.4, 11.0, 0.7, "What is next and how to engage",
         size=32, bold=True, color=CRSF_GREEN, font="Aptos Display")
add_text(s, 1.7, 1.15, 11.0, 0.5,
         "v2.0.1 and v2.1.0 are queued behind Cardinal compute; the manuscript is targeted at Forest Ecology and Management.",
         size=14, color=MED_GRAY)

# 4 milestones
milestones = [
    ("Now", "FEM submission preparation",
     "Final v0.10 polish, internal review, then submit. Outreach package staged at CRSF news + UMaine Comms."),
    ("Next 2 weeks", "v2.0.1 release",
     "Adds 30 m quantile-RF uncertainty raster for ESI v7. Cardinal job 11553490 is mid-array (4 of 40 running)."),
    ("Next 1–2 months", "v2.1.0 release",
     "30 m wall-to-wall surfaces (composite + components) once Zenodo quota expansion approved. NPP upgrade to MOD17 500 m baked in."),
    ("Ongoing", "Collaborator integrations",
     "FVS DG-Kuehne calibrations, CFRU growth-and-yield benchmarks, Carbon Field Pilot validation, ForCAST applications."),
]
y = 1.95
for i, (when, name, desc) in enumerate(milestones):
    add_rect(s, 1.0, y, 0.55, 0.55, CRSF_GREEN)
    add_text(s, 1.0, y+0.05, 0.55, 0.45, str(i+1), size=22, bold=True,
             color=WHITE, font="Aptos Display", align=PP_ALIGN.CENTER)
    add_text(s, 1.75, y+0.05, 2.2, 0.5, when, size=14, bold=True, color=CRSF_LIGHT,
             font="Aptos")
    add_text(s, 4.0, y+0.05, 8.5, 0.5, name, size=18, bold=True, color=CHARCOAL,
             font="Aptos Display")
    add_text(s, 4.0, y+0.55, 8.5, 0.65, desc, size=13, color=MED_GRAY)
    y += 1.2

add_text(s, 1.0, 6.6, 11.3, 0.35,
         "Engage: aaron.weiskittel@maine.edu  |  GitHub: holoros/fvs-conus  |  Zenodo concept 10.5281/zenodo.20515034",
         size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ----------------------------------------------------------------------------
# Slide 13: Closing
# ----------------------------------------------------------------------------
s = blank()
add_rect(s, 0, 0, 13.333, 7.5, CRSF_GREEN)
try:
    s.shapes.add_picture(CRSF_LOGO, Inches(11.6), Inches(0.4),
                         width=Inches(1.3), height=Inches(1.3))
except Exception:
    pass
add_text(s, 0.8, 2.4, 12.0, 1.5,
         "Questions, applications, collaborations?",
         size=44, bold=True, color=WHITE, font="Aptos Display",
         align=PP_ALIGN.CENTER)
add_text(s, 0.8, 4.4, 12.0, 0.6,
         "aaron.weiskittel@maine.edu",
         size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 5.2, 12.0, 0.5,
         "Center for Research on Sustainable Forests  |  University of Maine",
         size=18, color=CRSF_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 5.9, 12.0, 0.4,
         "github.com/holoros/fvs-conus  |  10.5281/zenodo.20515034  |  v0.10 manuscript in revision for FEM",
         size=14, color=CRSF_LIGHT, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Wrote: {OUT}")
print(f"Slides: {len(prs.slides)}")
